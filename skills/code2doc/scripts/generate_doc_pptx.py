#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PowerPoint Documentation Generator
코드 분석 결과를 PowerPoint 발표자료로 생성

Usage:
    python generate_doc_pptx.py --project "ProjectName" --output "/path/to/output"
    
Dependencies:
    pip install python-pptx --break-system-packages
"""

import os
from datetime import datetime
from typing import Dict, List, Optional, Tuple

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RgbColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("Warning: python-pptx not installed. Run: pip install python-pptx --break-system-packages")


class PowerPointDocGenerator:
    """PowerPoint 문서 생성기"""
    
    # 색상 팔레트
    COLORS = {
        'dark_blue': RgbColor(26, 54, 93),      # #1a365d
        'light_blue': RgbColor(49, 130, 206),   # #3182ce
        'white': RgbColor(255, 255, 255),
        'light_gray': RgbColor(247, 250, 252),  # #f7fafc
        'dark_gray': RgbColor(45, 55, 72),      # #2d3748
        'accent_green': RgbColor(56, 161, 105), # #38a169
        'accent_orange': RgbColor(221, 107, 32) # #dd6b20
    }
    
    def __init__(self, project_name: str, output_dir: str = "/mnt/user-data/outputs"):
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx is required. Install with: pip install python-pptx --break-system-packages")
            
        self.project_name = project_name
        self.output_dir = output_dir
        self.prs = Presentation()
        
        # 슬라이드 크기 설정 (16:9)
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        
    def _get_blank_slide(self):
        """빈 슬라이드 레이아웃 반환"""
        return self.prs.slide_layouts[6]  # Blank layout
        
    def _add_title_shape(self, slide, text: str, top: float = 0.5, 
                         font_size: int = 36, color=None):
        """제목 텍스트 박스 추가"""
        if color is None:
            color = self.COLORS['dark_blue']
            
        shape = slide.shapes.add_textbox(
            Inches(0.5), Inches(top),
            Inches(12.333), Inches(0.8)
        )
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = True
        p.font.color.rgb = color
        return shape
        
    def _add_body_text(self, slide, text: str, left: float = 0.5, 
                       top: float = 1.5, width: float = 12.333, 
                       height: float = 5.5, font_size: int = 18):
        """본문 텍스트 박스 추가"""
        shape = slide.shapes.add_textbox(
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = self.COLORS['dark_gray']
        return shape
        
    def _add_bullet_text(self, slide, items: List[str], left: float = 0.5,
                         top: float = 1.5, width: float = 12.333,
                         height: float = 5.5, font_size: int = 18):
        """글머리 기호 텍스트 추가"""
        shape = slide.shapes.add_textbox(
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )
        tf = shape.text_frame
        tf.word_wrap = True
        
        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(font_size)
            p.font.color.rgb = self.COLORS['dark_gray']
            p.space_after = Pt(12)
            
        return shape
        
    def _add_code_box(self, slide, code: str, left: float = 0.5,
                      top: float = 1.5, width: float = 12.333,
                      height: float = 5.5):
        """코드 박스 추가"""
        # 배경 사각형
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.COLORS['light_gray']
        bg.line.color.rgb = self.COLORS['light_blue']
        
        # 코드 텍스트
        shape = slide.shapes.add_textbox(
            Inches(left + 0.2), Inches(top + 0.2),
            Inches(width - 0.4), Inches(height - 0.4)
        )
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = code
        p.font.size = Pt(12)
        p.font.name = "Courier New"
        p.font.color.rgb = self.COLORS['dark_gray']
        
        return shape
        
    def _add_table(self, slide, headers: List[str], rows: List[List[str]],
                   left: float = 0.5, top: float = 1.5,
                   width: float = 12.333, row_height: float = 0.5):
        """테이블 추가"""
        num_rows = len(rows) + 1  # +1 for header
        num_cols = len(headers)
        col_width = width / num_cols
        
        table = slide.shapes.add_table(
            num_rows, num_cols,
            Inches(left), Inches(top),
            Inches(width), Inches(row_height * num_rows)
        ).table
        
        # 헤더 설정
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.COLORS['dark_blue']
            p = cell.text_frame.paragraphs[0]
            p.font.bold = True
            p.font.color.rgb = self.COLORS['white']
            p.font.size = Pt(14)
            
        # 데이터 행
        for row_idx, row in enumerate(rows, 1):
            for col_idx, value in enumerate(row):
                cell = table.cell(row_idx, col_idx)
                cell.text = str(value)
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(12)
                p.font.color.rgb = self.COLORS['dark_gray']
                
                # 줄무늬 배경
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self.COLORS['light_gray']
                    
        return table
        
    def add_title_slide(self, subtitle: str = "Technical Documentation"):
        """표지 슬라이드 추가"""
        slide = self.prs.slides.add_slide(self._get_blank_slide())
        
        # 배경 색상
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0,
            self.prs.slide_width, self.prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.COLORS['dark_blue']
        background.line.fill.background()
        
        # 메인 제목
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5),
            Inches(12.333), Inches(1.5)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = self.project_name
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = self.COLORS['white']
        p.alignment = PP_ALIGN.CENTER
        
        # 부제목
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4),
            Inches(12.333), Inches(0.8)
        )
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(28)
        p.font.color.rgb = self.COLORS['light_blue']
        p.alignment = PP_ALIGN.CENTER
        
        # 날짜
        date_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(6.5),
            Inches(12.333), Inches(0.5)
        )
        tf = date_box.text_frame
        p = tf.paragraphs[0]
        p.text = datetime.now().strftime("%Y-%m-%d")
        p.font.size = Pt(18)
        p.font.color.rgb = self.COLORS['white']
        p.alignment = PP_ALIGN.CENTER
        
        return slide
        
    def add_toc_slide(self, sections: List[str]):
        """목차 슬라이드 추가"""
        slide = self.prs.slides.add_slide(self._get_blank_slide())
        self._add_title_shape(slide, "Table of Contents")
        
        # 목차 항목
        items_text = "\n".join([f"{i+1}. {section}" for i, section in enumerate(sections)])
        
        shape = slide.shapes.add_textbox(
            Inches(1), Inches(1.5),
            Inches(11), Inches(5.5)
        )
        tf = shape.text_frame
        tf.word_wrap = True
        
        for i, section in enumerate(sections):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"{i+1}.  {section}"
            p.font.size = Pt(24)
            p.font.color.rgb = self.COLORS['dark_gray']
            p.space_after = Pt(20)
            
        return slide
        
    def add_content_slide(self, title: str, content: str = None,
                          bullets: List[str] = None, code: str = None):
        """일반 콘텐츠 슬라이드 추가"""
        slide = self.prs.slides.add_slide(self._get_blank_slide())
        self._add_title_shape(slide, title)
        
        if content:
            self._add_body_text(slide, content)
        elif bullets:
            self._add_bullet_text(slide, bullets)
        elif code:
            self._add_code_box(slide, code)
            
        return slide
        
    def add_table_slide(self, title: str, headers: List[str], 
                        rows: List[List[str]]):
        """테이블 슬라이드 추가"""
        slide = self.prs.slides.add_slide(self._get_blank_slide())
        self._add_title_shape(slide, title)
        self._add_table(slide, headers, rows)
        return slide
        
    def add_diagram_slide(self, title: str, diagram_text: str,
                          description: str = None):
        """다이어그램 슬라이드 추가 (ASCII 다이어그램)"""
        slide = self.prs.slides.add_slide(self._get_blank_slide())
        self._add_title_shape(slide, title)
        
        # 다이어그램 박스
        self._add_code_box(slide, diagram_text, top=1.5, height=4)
        
        # 설명 (선택)
        if description:
            self._add_body_text(slide, description, top=5.7, height=1.3, font_size=14)
            
        return slide
        
    def add_two_column_slide(self, title: str, 
                              left_title: str, left_items: List[str],
                              right_title: str, right_items: List[str]):
        """2열 슬라이드 추가"""
        slide = self.prs.slides.add_slide(self._get_blank_slide())
        self._add_title_shape(slide, title)
        
        # 왼쪽 열 제목
        left_title_shape = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.3),
            Inches(6), Inches(0.5)
        )
        tf = left_title_shape.text_frame
        p = tf.paragraphs[0]
        p.text = left_title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = self.COLORS['light_blue']
        
        # 왼쪽 열 내용
        self._add_bullet_text(slide, left_items, left=0.5, top=1.8, 
                              width=6, height=5, font_size=16)
        
        # 오른쪽 열 제목
        right_title_shape = slide.shapes.add_textbox(
            Inches(6.8), Inches(1.3),
            Inches(6), Inches(0.5)
        )
        tf = right_title_shape.text_frame
        p = tf.paragraphs[0]
        p.text = right_title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = self.COLORS['light_blue']
        
        # 오른쪽 열 내용
        self._add_bullet_text(slide, right_items, left=6.8, top=1.8,
                              width=6, height=5, font_size=16)
        
        return slide
        
    def add_qa_slide(self):
        """Q&A 슬라이드 추가"""
        slide = self.prs.slides.add_slide(self._get_blank_slide())
        
        # 배경 색상
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0,
            self.prs.slide_width, self.prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.COLORS['dark_blue']
        background.line.fill.background()
        
        # Q&A 텍스트
        qa_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(3),
            Inches(12.333), Inches(1.5)
        )
        tf = qa_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Q & A"
        p.font.size = Pt(72)
        p.font.bold = True
        p.font.color.rgb = self.COLORS['white']
        p.alignment = PP_ALIGN.CENTER
        
        # 감사 메시지
        thanks_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(5),
            Inches(12.333), Inches(0.8)
        )
        tf = thanks_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Thank you for your attention"
        p.font.size = Pt(24)
        p.font.color.rgb = self.COLORS['light_blue']
        p.alignment = PP_ALIGN.CENTER
        
        return slide
        
    def save(self, filename: Optional[str] = None) -> str:
        """프레젠테이션 저장"""
        if filename is None:
            filename = f"{self.project_name.lower()}_presentation.pptx"
            
        filepath = os.path.join(self.output_dir, filename)
        self.prs.save(filepath)
        return filepath


# =============================================================================
# 템플릿 함수들
# =============================================================================

def create_basic_presentation(project_name: str, 
                               output_dir: str = "/mnt/user-data/outputs") -> PowerPointDocGenerator:
    """기본 발표자료 템플릿 생성"""
    ppt = PowerPointDocGenerator(project_name, output_dir)
    
    # 표지
    ppt.add_title_slide()
    
    # 목차
    ppt.add_toc_slide([
        "Overview",
        "Architecture",
        "Workflow",
        "Key Components",
        "Build & Run"
    ])
    
    return ppt


def example_usage():
    """사용 예시"""
    ppt = create_basic_presentation("MyProject")
    
    # 개요
    ppt.add_content_slide(
        "1. Overview",
        bullets=[
            "Project purpose and goals",
            "Main features and capabilities",
            "Target users and use cases",
            "Technology stack overview"
        ]
    )
    
    # 아키텍처 다이어그램
    diagram = """
    +-------------+     +-------------+     +-------------+
    |   Client    |---->|   Server    |---->|  Database   |
    | (Frontend)  |     |   (API)     |     | (Storage)   |
    +-------------+     +-------------+     +-------------+
          |                   |
          v                   v
    +-------------+     +-------------+
    |    Auth     |     |   Cache     |
    +-------------+     +-------------+
    """
    ppt.add_diagram_slide("2. Architecture", diagram, 
                          "Layered architecture with caching and authentication")
    
    # 기술 스택 테이블
    ppt.add_table_slide(
        "Technology Stack",
        ["Category", "Technology", "Version"],
        [
            ["Language", "Python", "3.10+"],
            ["Framework", "FastAPI", "0.100+"],
            ["Database", "PostgreSQL", "15"],
            ["Cache", "Redis", "7.0"]
        ]
    )
    
    # 워크플로우
    ppt.add_content_slide(
        "3. Workflow",
        bullets=[
            "Step 1: Receive request from client",
            "Step 2: Validate and authenticate",
            "Step 3: Process business logic",
            "Step 4: Query/update database",
            "Step 5: Return response"
        ]
    )
    
    # 2열 레이아웃
    ppt.add_two_column_slide(
        "4. Key Components",
        "Frontend", ["React UI", "State management", "API client", "Routing"],
        "Backend", ["REST API", "Authentication", "Business logic", "Data layer"]
    )
    
    # 빌드 명령어
    build_code = """# Install dependencies
pip install -r requirements.txt

# Run development server
python main.py --dev

# Run tests
pytest tests/

# Build for production
python build.py --release"""
    ppt.add_content_slide("5. Build & Run", code=build_code)
    
    # Q&A
    ppt.add_qa_slide()
    
    # 저장
    filepath = ppt.save()
    print(f"Presentation saved: {filepath}")
    
    return ppt


if __name__ == "__main__":
    import argparse
    
    parser = argparse.ArgumentParser(description="Generate PowerPoint documentation")
    parser.add_argument("--project", default="Project", help="Project name")
    parser.add_argument("--output", default="/mnt/user-data/outputs", help="Output directory")
    parser.add_argument("--example", action="store_true", help="Generate example presentation")
    
    args = parser.parse_args()
    
    if args.example:
        example_usage()
    else:
        print(f"Creating template for: {args.project}")
        ppt = create_basic_presentation(args.project, args.output)
        
        # 기본 슬라이드들 추가
        ppt.add_content_slide("Overview", bullets=["Add your content here"])
        ppt.add_qa_slide()
        
        filepath = ppt.save()
        print(f"Template saved: {filepath}")
